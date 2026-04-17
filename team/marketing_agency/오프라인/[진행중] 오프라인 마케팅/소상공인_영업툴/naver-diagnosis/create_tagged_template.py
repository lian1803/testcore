"""
template.pptx → template.pptx (태그 버전으로 1회 변환)
placeholder 텍스트를 {{태그}} 형식으로 교체.
실행 후 template.pptx가 덮어쓰여짐.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from copy import deepcopy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

SRC = 'template.pptx'
DST = 'template.pptx'  # 덮어씀


# ── 헬퍼 ──────────────────────────────────────────────────────────────────────

def _para_text(para):
    return ''.join(r.text for r in para.runs)

def _replace_para(para, old, new):
    full = _para_text(para)
    if old not in full:
        return False
    replaced = full.replace(old, new)
    if para.runs:
        para.runs[0].text = replaced
        for r in para.runs[1:]:
            r.text = ''
    return True

def _replace_shape(shape, old, new):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        _replace_para(para, old, new)

def _replace_recursive(shape, old, new):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            _replace_recursive(s, old, new)
    else:
        _replace_shape(shape, old, new)

def _replace_slide(slide, old, new):
    for s in slide.shapes:
        _replace_recursive(s, old, new)

def _collect_text_shapes(shape, result):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            _collect_text_shapes(s, result)
    elif shape.has_text_frame:
        result.append(shape)

def _set_shape_content(shape, text):
    """shape 전체 내용을 단일 텍스트로 교체."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for i, para in enumerate(tf.paragraphs):
        if i == 0:
            if para.runs:
                para.runs[0].text = text
                for r in para.runs[1:]:
                    r.text = ''
        else:
            for r in para.runs:
                r.text = ''

def _append_paragraph(tf, text):
    last_p = tf.paragraphs[-1]._p
    new_p = deepcopy(last_p)
    runs = new_p.findall(qn('a:r'))
    if runs:
        runs[0].find(qn('a:t')).text = text
        for r in runs[1:]:
            new_p.remove(r)
    tf._txBody.append(new_p)


# ── 슬라이드별 태그 삽입 ──────────────────────────────────────────────────────

prs = Presentation(SRC)
slides = prs.slides

# ── 슬라이드 1: 커버 ──
# "분석 보고서" 단락 아래에 {{업체명}} 단락 추가
for shape in slides[0].shapes:
    if shape.has_text_frame and '보고서' in shape.text_frame.text:
        _append_paragraph(shape.text_frame, '{{업체명}}')
        break

# ── 슬라이드 2: 업체 분석 ──
s2_map = {
    ' (업체명) 분석':    '{{업체명}} 분석',
    '(업체명) 분석':     '{{업체명}} 분석',
    '한의원':           '{{업종}}',
    '59위':             '{{순위}}',
    '집중패키지':        '{{추천패키지}}',
    '가망 고객 분석':    '{{가망고객분석}}',
    'seo 최적화':        '{{주요문제}}',
    'Pc : 80회':         '{{PC검색량}}',
    '모바일 : 80회':     '{{모바일검색량}}',
    '전체 : 80회':       '{{전체검색량}}',
}
for old, new in s2_map.items():
    _replace_slide(slides[1], old, new)

# "Hours a week..." — 그룹 안에도 있음
_replace_slide(slides[1], 'Hours a week spent doing automable jobs', '{{주소}}')

# ── 슬라이드 3: 키워드 테이블 ──
_replace_slide(slides[2], '(업체명) 노출전략', '{{업체명}} 노출전략')

kw_shapes  = [s for s in slides[2].shapes
              if s.has_text_frame and '디자인 시벌' in s.text_frame.text]
vol_shapes = [s for s in slides[2].shapes
              if s.has_text_frame and '20,180' in s.text_frame.text]

for i, shape in enumerate(kw_shapes):
    _replace_shape(shape, '디자인 시벌', '{{키워드%d}}' % (i + 1))
for i, shape in enumerate(vol_shapes):
    _replace_shape(shape, '20,180', '{{검색량%d}}' % (i + 1))

# ── 슬라이드 4: 검색 결과 ──
for shape in slides[3].shapes:
    if shape.has_text_frame and '(업체명)' in shape.text_frame.text \
            and '검색' not in shape.text_frame.text:
        _replace_shape(shape, '(업체명)', '{{업체명}}')
        break

# ── 슬라이드 5: 종합 진단 — description 자리에 태그 ──
desc_tag_map = {
    '방문결정에':          '{{진단_대표사진}}',
    '소개 글':             '{{진단_소개오시는길}}',
    '현재 알림':           '{{진단_알림이벤트}}',
    '네이버는 행정구역':    '{{진단_대표키워드}}',
    '세팅이 되어 있으나':   '{{진단_메뉴세팅}}',
    '예약 , 톡톡':          '{{진단_예약세팅}}',
    '영수증리뷰':           '{{진단_영수증리뷰}}',
    '블로그가 있으나':      '{{진단_블로그노출}}',
    '플레이스에 링크 삽입이':'{{진단_지역커뮤니티}}',
    '외부채널 링크':        '{{진단_외부채널}}',
}
all_s5 = []
for s in slides[4].shapes:
    _collect_text_shapes(s, all_s5)

for anchor, tag in desc_tag_map.items():
    for shape in all_s5:
        if anchor in shape.text_frame.text:
            _set_shape_content(shape, tag)
            break

# ── 슬라이드 8: 제안서 ──
_replace_slide(slides[7], '상호명', '{{상호명}}')
_replace_slide(slides[7], '주목 패키지', '{{패키지}}')
_replace_slide(slides[7], '주목패키지', '{{패키지}}')
# 연쇄 방지: 정상가 먼저 토큰으로, 나중에 태그로
_replace_slide(slides[7], '850,000', '__LIST_PRICE__')
_replace_slide(slides[7], '560,000', '{{약정가}}')
_replace_slide(slides[7], '__LIST_PRICE__', '{{정상가}}')
_replace_slide(slides[7], '3,360,000', '{{합계6개월}}')

prs.save(DST)
print("✅ 태그 템플릿 생성 완료:", DST)

# 검증: 삽입된 태그 목록 출력
print("\n─ 삽입된 태그 확인 ─")
prs2 = Presentation(DST)
found_tags = set()
import re
for slide in prs2.slides:
    all_shapes = []
    for s in slide.shapes:
        _collect_text_shapes(s, all_shapes)
    for shape in all_shapes:
        for m in re.findall(r'\{\{[^}]+\}\}', shape.text_frame.text):
            found_tags.add(m)

for tag in sorted(found_tags):
    print(' ', tag)
