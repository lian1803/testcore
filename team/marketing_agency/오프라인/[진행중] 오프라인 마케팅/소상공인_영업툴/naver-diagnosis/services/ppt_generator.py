"""
PPT 제안서 생성기 — 태그 기반 (v2)
template.pptx의 {{태그}}를 실데이터로 교체하여 제안서 생성.

태그 목록: create_tagged_template.py 참고
"""
import json
import os
from datetime import datetime
from typing import Dict, Any, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from services.prescription_generator import PrescriptionGenerator

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'template.pptx')

PACKAGE_BY_GRADE = {
    'D': '시선패키지',
    'C': '집중패키지',
    'B': '주목패키지',
    'A': '주목패키지',
}

PACKAGE_PRICES = {
    '시선패키지': {'정상가': '560,000', '약정가': '380,000', '합계6개월': '2,280,000'},  # 시선 38만원/월 (앵커 56만)
    '집중패키지': {'정상가': '850,000', '약정가': '560,000', '합계6개월': '3,360,000'},  # 집중 56만원/월 (앵커 85만)
    '주목패키지': {'정상가': '1,270,000', '약정가': '950,000', '합계6개월': '5,700,000'},  # 주목 95만원/월 (앵커 127만)
}

CATEGORY_LABELS = {
    'photo': '대표사진',
    'review': '리뷰 관리',
    'blog': '블로그 리뷰',
    'info': '기본정보 완성',
    'keyword': '키워드 세팅',
    'convenience': '알림/이벤트',
    'engagement': '고객 소통',
}


# ── 태그 교체 엔진 ────────────────────────────────────────────────────────────

def _replace_para(para, tag: str, value: str):
    full = ''.join(r.text for r in para.runs)
    if tag not in full:
        return
    replaced = full.replace(tag, value)
    if para.runs:
        para.runs[0].text = replaced
        for r in para.runs[1:]:
            r.text = ''


def _replace_in_shape(shape, tag: str, value: str, font_size_pt: int = None):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        _replace_para(para, tag, value)
    if font_size_pt:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.font.size = Pt(font_size_pt)


def _replace_all(prs: Presentation, tags: Dict[str, str]):
    """프레젠테이션 전체에서 {{태그}} → 값 교체 (그룹 재귀 포함)."""
    def _walk(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                _walk(s)
        elif shape.has_text_frame:
            for tag, value in tags.items():
                _replace_in_shape(shape, tag, value)

    for slide in prs.slides:
        for shape in slide.shapes:
            _walk(shape)


def _replace_title_fit(prs: Presentation, tag: str, value: str, font_size_pt: int = 36):
    """타이틀 shape에서 태그 교체 + 폰트 크기 축소 (텍스트 넘침 방지)."""
    def _walk(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                _walk(s)
        elif shape.has_text_frame and tag in shape.text_frame.text:
            _replace_in_shape(shape, tag, value, font_size_pt=font_size_pt)

    for slide in prs.slides:
        for shape in slide.shapes:
            _walk(shape)


# ── 데이터 → 태그 dict 생성 ───────────────────────────────────────────────────

class PPTGenerator:

    def __init__(self, output_dir: str = "ppt_output"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def _parse(self, value, default=None):
        if default is None:
            default = []
        if not value:
            return default
        if isinstance(value, (list, dict)):
            return value
        try:
            return json.loads(value)
        except Exception:
            return default

    def _build_tags(self, data: Dict) -> Dict[str, str]:
        """data dict → {{태그}}: 값 매핑 생성"""

        # 기본 필드
        name      = data.get('business_name') or '업체명'
        category  = data.get('category') or '소상공인'
        rank      = data.get('naver_place_rank') or 0
        address   = data.get('address') or '주소 미확인'
        grade     = data.get('grade') or 'D'
        lost      = data.get('estimated_lost_customers') or 0

        keywords  = self._parse(data.get('keywords'))
        if not keywords:
            keywords = self._parse(data.get('related_keywords'))
        improvement_points = self._parse(data.get('improvement_points'))

        package   = PACKAGE_BY_GRADE.get(grade, '시선패키지')
        prices    = PACKAGE_PRICES[package]

        # 주요 문제 레이블
        top_issue = '개선 필요'
        if improvement_points:
            top_issue = CATEGORY_LABELS.get(improvement_points[0].get('category', ''), '개선 필요')

        # 키워드 검색량
        total_search = sum(k.get('search_volume', 0) for k in keywords)
        pc_search    = int(total_search * 0.3)
        mobile_search = int(total_search * 0.7)

        # 가망 고객 분석 문구
        lost_txt = (f'⚠ 매달 약 {lost}명이 경쟁사로 이탈 중'
                    if lost > 0 else '가망 고객 분석')

        # 처방전 생성
        rx = PrescriptionGenerator().generate(
            grade=grade,
            score=data.get('total_score', 0),
            weak_items=improvement_points,
            business_name=name,
            category=category,
            data=data,
        )

        tags: Dict[str, str] = {
            # 공통
            '{{업체명}}':       name,
            '{{업종}}':         category,
            '{{순위}}':         f'{rank}위' if rank else '미확인',
            '{{주소}}':         address,
            '{{추천패키지}}':    package,
            '{{가망고객분석}}':  lost_txt,
            '{{주요문제}}':     top_issue,

            # 검색량
            '{{PC검색량}}':      f'Pc : {pc_search:,}회',
            '{{모바일검색량}}':  f'모바일 : {mobile_search:,}회',
            '{{전체검색량}}':    f'전체 : {total_search:,}회',

            # 제안서 (가격)
            '{{상호명}}':       name,
            '{{패키지}}':       package,
            '{{정상가}}':       prices['정상가'],   # 앵커 가격 (취소선)
            '{{약정가}}':       prices['약정가'],   # 실제 판매가
            '{{합계6개월}}':    prices['합계6개월'],

            # 처방전
            '{{처방요약}}':     rx['prescription_text'],
            '{{손익분기}}':     rx['breakeven']['summary'],
        }

        # 키워드 테이블 (7행)
        for i in range(1, 8):
            kw = keywords[i - 1] if i - 1 < len(keywords) else {}
            tags[f'{{{{키워드{i}}}}}']  = kw.get('keyword', '-') or '-'
            vol = kw.get('search_volume', 0)
            tags[f'{{{{검색량{i}}}}}']  = f'{vol:,}' if vol else '-'

        # 이달 처방 우선순위 (최대 3개)
        for i in range(1, 4):
            p_list = rx.get('priorities', [])
            if i - 1 < len(p_list):
                p = p_list[i - 1]
                tags[f'{{{{처방{i}순위}}}}'] = f"{p['action']} — {p['why']}"
            else:
                tags[f'{{{{처방{i}순위}}}}'] = '-'

        # 슬라이드 5 진단 텍스트
        diag = self._build_diagnosis_texts(data, keywords)
        tags.update(diag)

        return tags

    def _build_diagnosis_texts(self, data: Dict, keywords: List) -> Dict[str, str]:
        """슬라이드 5 진단 항목별 텍스트 생성"""
        photo_count    = data.get('photo_count', 0) or 0
        visitor_review = data.get('visitor_review_count', 0) or 0
        receipt_review = data.get('receipt_review_count', 0) or 0
        blog_count     = data.get('blog_review_count', 0) or 0
        menu_count     = data.get('menu_count', 0) or 0
        has_intro      = bool(data.get('has_intro'))
        has_directions = bool(data.get('has_directions'))
        has_booking    = bool(data.get('has_booking'))
        has_talktalk   = bool(data.get('has_talktalk'))
        has_smartcall  = bool(data.get('has_smartcall'))
        has_news       = bool(data.get('has_news'))
        has_coupon     = bool(data.get('has_coupon'))
        has_instagram  = bool(data.get('has_instagram'))
        has_kakao      = bool(data.get('has_kakao'))
        has_price      = bool(data.get('has_price'))
        has_menu       = bool(data.get('has_menu'))
        has_menu_desc  = bool(data.get('has_menu_description'))
        kw_count       = len(keywords)
        total_search   = sum(k.get('search_volume', 0) for k in keywords)

        def _missing(*items):
            return '·'.join(items)

        # 대표사진
        if photo_count == 0:
            photo = '사진 0장 — 방문 결정의 첫 번째 요소가 비어있어요.'
        elif photo_count < 10:
            photo = f'사진 {photo_count}장 — 이 지역 상위 업체 평균 대비 미달이에요.'
        else:
            photo = f'사진 {photo_count}장 — 경쟁 업체 대비 부족한 상태예요.'

        # 소개/오시는길
        m = [x for x, v in [('소개글', not has_intro), ('오시는길', not has_directions)] if v]
        info = (f'{_missing(*m)} 미입력 — 검색 키워드 노출 및 신뢰도에 영향을 미칩니다.'
                if m else '소개글·오시는길 입력됨 — 하지만 키워드 노출 경쟁력을 확인해야 해요.')

        # 알림/이벤트
        m = [x for x, v in [('새소식', not has_news), ('쿠폰', not has_coupon)] if v]
        event = (f'{_missing(*m)} 미설정 — 플레이스 활성화 지수에 영향을 줍니다.'
                 if m else '알림·이벤트 운영 중 — 업데이트 주기와 활성화 지수를 점검해야 해요.')

        # 대표키워드
        kw = (f'키워드 {kw_count}개 · 월 검색량 {total_search:,}회 — 더 많은 검색 유입이 가능한 상태예요.'
              if kw_count else '키워드 미설정 — 검색 노출이 어렵습니다. 핵심 지역 키워드 세팅 필요합니다.')

        # 상품(메뉴)
        if not has_menu or menu_count == 0:
            menu = '메뉴 미등록 — 방문 전 정보를 찾는 손님들이 이탈하고 있어요.'
        elif not has_price:
            menu = f'메뉴 {menu_count}개 · 가격 미입력 — 가격을 확인하러 경쟁사로 넘어가고 있어요.'
        elif not has_menu_desc:
            menu = f'메뉴 {menu_count}개 · 가격 입력됨 — 상세 설명 보완을 권장합니다.'
        else:
            menu = f'메뉴 {menu_count}개 · 가격·설명 입력됨 — 최신성 유지 여부 점검 필요해요.'

        # 예약
        m = [x for x, v in [('예약', not has_booking), ('톡톡', not has_talktalk),
                              ('스마트콜', not has_smartcall)] if v]
        booking = (f'{_missing(*m)} 미설정 — 방문 유도 기회를 놓치고 있습니다.'
                   if m else '예약·톡톡·스마트콜 설정됨 — 고객 유입 채널 연결 상태예요.')

        # 영수증 리뷰
        receipt = (f'영수증 리뷰 0개 · 방문자 리뷰 {visitor_review}개 — 영수증 리뷰 유도가 필요합니다.'
                   if receipt_review == 0 else
                   f'영수증 리뷰 {receipt_review}개 · 방문자 리뷰 {visitor_review}개 — 지속 관리를 권장합니다.')

        # 블로그
        if blog_count == 0:
            blog = '블로그 리뷰 없음 — 네이버 검색에서 이 업체가 노출될 이유가 없어요.'
        elif blog_count < 5:
            blog = f'블로그 리뷰 {blog_count}개 — 경쟁사와 비교했을 때 검색 노출에서 밀리고 있어요.'
        else:
            blog = f'블로그 리뷰 {blog_count}개 — 최신성 유지 여부를 확인해야 해요.'

        # 커뮤니티
        community = ('카카오채널 연결됨 — 지역 커뮤니티 도달 범위를 확인해야 해요.'
                     if has_kakao else
                     '카카오채널 미연결 — 단골 관리 채널 부재. 개설 및 플레이스 연결이 필요합니다.')

        # 외부채널
        sns = ('인스타그램 연결됨 — SNS 유입 현황을 확인해야 해요.'
               if has_instagram else
               '인스타그램 미연결 — SNS 유입 경로 없음. 채널 개설 및 플레이스 연결 필요합니다.')

        return {
            '{{진단_대표사진}}':     photo,
            '{{진단_소개오시는길}}': info,
            '{{진단_알림이벤트}}':   event,
            '{{진단_대표키워드}}':   kw,
            '{{진단_메뉴세팅}}':     menu,
            '{{진단_예약세팅}}':     booking,
            '{{진단_영수증리뷰}}':   receipt,
            '{{진단_블로그노출}}':   blog,
            '{{진단_지역커뮤니티}}': community,
            '{{진단_외부채널}}':     sns,
        }

    # ── generate ─────────────────────────────────────────────────────────────

    def generate(self, data: Dict[str, Any]) -> str:
        prs = Presentation(TEMPLATE_PATH)

        tags = self._build_tags(data)

        # 타이틀 shape (슬라이드 3·4) — 폰트 크기 축소 적용
        # 나머지는 일반 교체
        title_tags = {'{{업체명}} 노출전략', '{{업체명}}'}

        # 1차: 타이틀 태그 (폰트 축소)
        _replace_title_fit(prs, '{{업체명}} 노출전략',
                           tags['{{업체명}}'] + ' 노출전략', font_size_pt=36)
        _replace_title_fit(prs, '{{업체명}}',
                           tags['{{업체명}}'], font_size_pt=36)

        # 2차: 나머지 태그 전체 교체 (업체명 태그는 이미 처리됨이므로 제외)
        remaining = {k: v for k, v in tags.items() if k != '{{업체명}}'}
        _replace_all(prs, remaining)

        # 저장
        safe_name = tags['{{업체명}}'].replace('/', '_').replace('\\', '_').replace(' ', '_')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{safe_name}_{timestamp}.pptx"
        prs.save(os.path.join(self.output_dir, filename))
        return filename
